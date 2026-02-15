"""
Fluir - Servico de Exportacao (Excel + PPT)
"""

# Versao do formato PPT (incrementar ao alterar layout/graficos/copywriting)
# Usado pelo endpoint /api/version para diagnostico (backend novo vs antigo)
PPT_FORMAT_VERSION = "2.0"

import io
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt, Inches

from copsoq_data import DIMENSIONS

from ppt_copy import (
    FECHAMENTO_CTA,
    INTRO_RECOMENDACOES,
    RESUMO_EXECUTIVO,
    ROI_BULLETS,
    ROI_TITULO,
)


# ──────── Color helpers (Excel) ────────
FILL_GREEN = PatternFill(start_color="D5F5E3", end_color="D5F5E3", fill_type="solid")
FILL_YELLOW = PatternFill(start_color="FEF9E7", end_color="FEF9E7", fill_type="solid")
FILL_RED = PatternFill(start_color="FADBD8", end_color="FADBD8", fill_type="solid")
FILL_HEADER = PatternFill(start_color="0F4C75", end_color="0F4C75", fill_type="solid")
FILL_SUBHEADER = PatternFill(start_color="3282B8", end_color="3282B8", fill_type="solid")
WHITE_FONT = Font(name="Inter", size=11, bold=True, color="FFFFFF")
NORMAL_FONT = Font(name="Inter", size=10, color="2C3E50")
BOLD_FONT = Font(name="Inter", size=10, bold=True, color="2C3E50")
THIN_BORDER = Border(
    left=Side(style="thin", color="CCCCCC"),
    right=Side(style="thin", color="CCCCCC"),
    top=Side(style="thin", color="CCCCCC"),
    bottom=Side(style="thin", color="CCCCCC"),
)

STATUS_FILL = {"green": FILL_GREEN, "yellow": FILL_YELLOW, "red": FILL_RED}
STATUS_LABEL = {"green": "Favorável", "yellow": "Atenção", "red": "Crítico"}
PRIORITY_LABEL = {"imediata": "Ação Imediata", "curto": "Curto Prazo", "medio": "Médio Prazo"}

# Formato 16:9 (widescreen)
SLIDE_WIDTH_16_9 = Inches(13.333)
SLIDE_HEIGHT_16_9 = Inches(7.5)
LOGO_SIZE_INCHES = 1.0
IMG_DIR = Path(__file__).resolve().parent / "static" / "img"

# Paleta PPT (branding Fluir)
COLOR_PRIMARY = RGBColor(0x0F, 0x4C, 0x75)   # Azul principal
COLOR_SECONDARY = RGBColor(0x32, 0x82, 0xB8)  # Azul secundario
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_YELLOW = RGBColor(0xF3, 0x9C, 0x12)
COLOR_RED = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GREEN_LIGHT = RGBColor(0xD5, 0xF5, 0xE3)
COLOR_YELLOW_LIGHT = RGBColor(0xFE, 0xF9, 0xE7)
COLOR_RED_LIGHT = RGBColor(0xFA, 0xDB, 0xD8)
STATUS_RGB = {
    "green": (COLOR_GREEN_LIGHT, COLOR_GREEN),
    "yellow": (COLOR_YELLOW_LIGHT, COLOR_YELLOW),
    "red": (COLOR_RED_LIGHT, COLOR_RED),
}

# Limites de truncamento para export
MAX_DIM_DESC_LEN = 50
MAX_REC_DESC_LEN = 80
MAX_PROSE_LEN = 2000


# ══════════════════════════════════════════════
# EXCEL EXPORT
# ══════════════════════════════════════════════

def export_excel(
    survey: Dict[str, Any],
    respondents_data: List[Dict[str, Any]],
    dim_scores_agg: List[Dict[str, Any]],
    kpis: Dict[str, Any],
    summary: Dict[str, Any],
    recommendations: List[Dict[str, Any]],
) -> io.BytesIO:
    """Gera Excel completo e retorna BytesIO."""
    wb = Workbook()
    CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
    LEFT = Alignment(horizontal="left", vertical="center", wrap_text=True)

    # ─── Aba 1: Resumo Executivo ───
    ws = wb.active
    ws.title = "Resumo Executivo"
    ws.sheet_properties.tabColor = "0F4C75"
    ws.column_dimensions["A"].width = 5
    ws.column_dimensions["B"].width = 40
    ws.column_dimensions["C"].width = 15
    ws.column_dimensions["D"].width = 15
    ws.column_dimensions["E"].width = 15
    ws.column_dimensions["F"].width = 40

    # Header
    ws.merge_cells("A1:F1")
    c = ws["A1"]
    c.value = f"Fluir — Relatório de Diagnóstico Psicossocial"
    c.font = Font(name="Inter", size=18, bold=True, color="FFFFFF")
    c.fill = FILL_HEADER
    c.alignment = CENTER
    ws.row_dimensions[1].height = 50

    ws.merge_cells("A2:F2")
    c = ws["A2"]
    c.value = f"{survey.get('company_name', '')} | {datetime.now().strftime('%d/%m/%Y')} | {summary.get('total', 0)} respondentes"
    c.font = Font(name="Inter", size=11, color="FFFFFF")
    c.fill = FILL_SUBHEADER
    c.alignment = CENTER
    ws.row_dimensions[2].height = 30

    # KPIs
    row = 4
    ws.merge_cells(f"A{row}:F{row}")
    ws[f"A{row}"].value = "INDICADORES-CHAVE (KPIs)"
    ws[f"A{row}"].font = Font(name="Inter", size=13, bold=True, color="FFFFFF")
    ws[f"A{row}"].fill = FILL_HEADER
    ws[f"A{row}"].alignment = CENTER

    row = 5
    for key, kpi in kpis.items():
        row += 1
        ws[f"B{row}"] = kpi["label"]
        ws[f"B{row}"].font = BOLD_FONT
        ws[f"C{row}"] = kpi["value"]
        ws[f"C{row}"].font = BOLD_FONT
        ws[f"C{row}"].alignment = CENTER
        ws[f"C{row}"].number_format = "0.00"
        ws[f"D{row}"] = kpi["status"]
        ws[f"D{row}"].alignment = CENTER
        ws[f"D{row}"].fill = STATUS_FILL.get(kpi["color"], FILL_YELLOW)
        for col in "BCDEF":
            ws[f"{col}{row}"].border = THIN_BORDER

    # Summary counts
    row += 2
    ws.merge_cells(f"A{row}:F{row}")
    ws[f"A{row}"].value = "RESUMO GERAL"
    ws[f"A{row}"].font = Font(name="Inter", size=13, bold=True, color="FFFFFF")
    ws[f"A{row}"].fill = FILL_HEADER
    ws[f"A{row}"].alignment = CENTER

    row += 1
    for label, count, color in [
        ("Dimensões Favoráveis", summary["green"], "27AE60"),
        ("Dimensões em Atenção", summary["yellow"], "F39C12"),
        ("Dimensões Críticas", summary["red"], "E74C3C"),
    ]:
        ws[f"B{row}"] = label
        ws[f"B{row}"].font = Font(name="Inter", size=11, bold=True, color=color)
        ws[f"C{row}"] = count
        ws[f"C{row}"].font = Font(name="Inter", size=14, bold=True, color=color)
        ws[f"C{row}"].alignment = CENTER
        row += 1

    # ─── Aba 2: Dimensões ───
    ws2 = wb.create_sheet("Dimensões")
    ws2.sheet_properties.tabColor = "3282B8"
    ws2.column_dimensions["A"].width = 5
    ws2.column_dimensions["B"].width = 35
    ws2.column_dimensions["C"].width = 12
    ws2.column_dimensions["D"].width = 14
    ws2.column_dimensions["E"].width = 12
    ws2.column_dimensions["F"].width = 28
    ws2.column_dimensions["G"].width = 35

    headers = ["#", "Dimensão", "Score", "Status", "Tipo", "Categoria", "Descrição"]
    for i, h in enumerate(headers, 1):
        c = ws2.cell(row=1, column=i, value=h)
        c.font = WHITE_FONT
        c.fill = FILL_HEADER
        c.alignment = CENTER

    for idx, d in enumerate(dim_scores_agg, 1):
        r = idx + 1
        ws2.cell(row=r, column=1, value=idx).alignment = CENTER
        ws2.cell(row=r, column=2, value=d["name"]).font = NORMAL_FONT
        ws2.cell(row=r, column=3, value=d["score"]).alignment = CENTER
        ws2.cell(row=r, column=3).number_format = "0.00"
        ws2.cell(row=r, column=4, value=STATUS_LABEL.get(d["status"], "")).alignment = CENTER
        ws2.cell(row=r, column=4).fill = STATUS_FILL.get(d["status"], FILL_YELLOW)
        ws2.cell(row=r, column=5, value="Risco" if d["type"] == "risk" else "Recurso").alignment = CENTER
        ws2.cell(row=r, column=6, value=d["category"])
        ws2.cell(row=r, column=7, value=d["description"])
        for col in range(1, 8):
            ws2.cell(row=r, column=col).border = THIN_BORDER

    # ─── Aba 3: Respostas Individuais ───
    ws3 = wb.create_sheet("Respostas Individuais")
    ws3.sheet_properties.tabColor = "27AE60"

    ws3.cell(row=1, column=1, value="ID").font = WHITE_FONT
    ws3.cell(row=1, column=1).fill = FILL_HEADER
    ws3.cell(row=1, column=1).alignment = CENTER
    ws3.column_dimensions["A"].width = 12

    dim_ids = list(DIMENSIONS.keys())
    for i, dim_id in enumerate(dim_ids, 2):
        c = ws3.cell(row=1, column=i, value=DIMENSIONS[dim_id]["name"])
        c.font = Font(name="Inter", size=9, bold=True, color="FFFFFF")
        c.fill = FILL_HEADER
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True, text_rotation=90)
        ws3.column_dimensions[get_column_letter(i)].width = 6

    for r_idx, resp in enumerate(respondents_data, 2):
        ws3.cell(row=r_idx, column=1, value=resp["display_id"]).alignment = CENTER
        for i, dim_id in enumerate(dim_ids, 2):
            score = resp.get("scores", {}).get(dim_id)
            if score is not None:
                cell = ws3.cell(row=r_idx, column=i, value=score)
                cell.alignment = CENTER
                cell.number_format = "0.00"
                status = resp.get("statuses", {}).get(dim_id, "yellow")
                cell.fill = STATUS_FILL.get(status, FILL_YELLOW)
                cell.border = THIN_BORDER

    # ─── Aba 4: Recomendações ───
    ws4 = wb.create_sheet("Recomendações")
    ws4.sheet_properties.tabColor = "7C83FD"
    ws4.column_dimensions["A"].width = 5
    ws4.column_dimensions["B"].width = 18
    ws4.column_dimensions["C"].width = 45
    ws4.column_dimensions["D"].width = 60

    for i, h in enumerate(["#", "Prioridade", "Título", "Descrição"], 1):
        c = ws4.cell(row=1, column=i, value=h)
        c.font = WHITE_FONT
        c.fill = FILL_HEADER
        c.alignment = CENTER

    for idx, rec in enumerate(recommendations, 1):
        r = idx + 1
        ws4.cell(row=r, column=1, value=idx).alignment = CENTER
        ws4.cell(row=r, column=2, value=PRIORITY_LABEL.get(rec.get("priority", ""), rec.get("priority", ""))).alignment = CENTER
        ws4.cell(row=r, column=3, value=rec.get("title", ""))
        ws4.cell(row=r, column=4, value=rec.get("description", ""))
        for col in range(1, 5):
            ws4.cell(row=r, column=col).border = THIN_BORDER
            ws4.cell(row=r, column=col).alignment = Alignment(vertical="center", wrap_text=True)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════
# PPT EXPORT
# ══════════════════════════════════════════════

def _build_category_scores(dim_scores_agg: List[Dict[str, Any]]) -> Dict[str, Dict[str, Any]]:
    """Agrupa dimensoes por categoria e calcula media. Retorna {nome_cat: {avg, type, status}}."""
    from collections import defaultdict
    cat_map: Dict[str, Dict[str, Any]] = defaultdict(lambda: {"sum": 0.0, "count": 0, "type": "resource"})
    for d in dim_scores_agg:
        cat = d.get("category", "Outros")
        cat_map[cat]["sum"] += d.get("score", 0)
        cat_map[cat]["count"] += 1
        cat_map[cat]["type"] = d.get("type", "resource")
    LOWER, UPPER = 2.33, 3.66
    result = {}
    for cat, data in cat_map.items():
        avg = data["sum"] / data["count"] if data["count"] else 0
        t = data["type"]
        if t == "risk":
            status = "green" if avg < LOWER else "red" if avg > UPPER else "yellow"
        else:
            status = "green" if avg > UPPER else "red" if avg < LOWER else "yellow"
        result[cat] = {"avg": avg, "type": t, "status": status}
    return result


def _render_radar_chart(category_scores: Dict[str, Dict[str, Any]]) -> io.BytesIO:
    """Gera grafico radar (8 categorias) como PNG em BytesIO."""
    import math
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    labels = list(category_scores.keys())
    values = [category_scores[l]["avg"] for l in labels]
    colors_map = {"green": "#6B9F7E", "yellow": "#D4A843", "red": "#C46B6B"}
    point_colors = [colors_map.get(category_scores[l]["status"], "#D4A843") for l in labels]
    num = len(labels)
    angles = [2 * math.pi * i / num for i in range(num)]
    values += [values[0]]
    angles += [angles[0]]
    fig, ax = plt.subplots(figsize=(8, 6), subplot_kw=dict(projection="polar"))
    ax.plot(angles, values, "o-", linewidth=2, color="#6B82A8")
    ax.fill(angles, values, alpha=0.15, color="#6B82A8")
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, size=9, wrap=True)
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _render_bar_chart(dim_scores_agg: List[Dict[str, Any]]) -> io.BytesIO:
    """Gera grafico de barras horizontais (26 dimensoes) como PNG em BytesIO."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    colors_map = {"green": "#6B9F7E", "yellow": "#D4A843", "red": "#C46B6B"}
    labels = [d.get("name", "")[:20] for d in dim_scores_agg]
    values = [d.get("score", 0) for d in dim_scores_agg]
    colors = [colors_map.get(d.get("status", "yellow"), "#D4A843") for d in dim_scores_agg]
    fig, ax = plt.subplots(figsize=(12, 8))
    y_pos = range(len(labels))
    ax.barh(y_pos, values, color=colors, height=0.7)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=9)
    ax.set_xlim(0, 5)
    ax.set_xlabel("Score")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_kpi_cards_slide(prs: Presentation, kpis: Dict[str, Any]) -> None:
    """Adiciona slide com KPIs em layout de cards 2x2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = "Indicadores-Chave (KPIs)"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    kpi_list = list(kpis.values())
    positions = [(0.5, 1.2), (7.0, 1.2), (0.5, 3.8), (7.0, 3.8)]
    card_w, card_h = Inches(6), Inches(2.2)
    for i, (left_in, top_in) in enumerate(positions):
        if i >= len(kpi_list):
            break
        kpi = kpi_list[i]
        left, top = Inches(left_in), Inches(top_in)
        color_key = kpi.get("color", "yellow")
        fill_rgb, _ = STATUS_RGB.get(color_key, STATUS_RGB["yellow"])
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
            shape.line.color.rgb = fill_rgb
        except Exception:
            shape = None
        box_left = left + Inches(0.3)
        box_top = top + Inches(0.3)
        val_box = slide.shapes.add_textbox(box_left, box_top, card_w - Inches(0.6), Inches(0.8))
        val_box.text_frame.paragraphs[0].text = f"{kpi['value']:.2f}"
        val_box.text_frame.paragraphs[0].font.size = Pt(36)
        val_box.text_frame.paragraphs[0].font.bold = True
        val_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        lbl_box = slide.shapes.add_textbox(box_left, top + Inches(1.1), card_w - Inches(0.6), Inches(0.5))
        lbl_box.text_frame.paragraphs[0].text = kpi.get("label", "")
        lbl_box.text_frame.paragraphs[0].font.size = Pt(12)
        lbl_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        status_box = slide.shapes.add_textbox(box_left, top + Inches(1.6), card_w - Inches(0.6), Inches(0.4))
        status_box.text_frame.paragraphs[0].text = kpi.get("status", "")
        status_box.text_frame.paragraphs[0].font.size = Pt(11)
        status_box.text_frame.paragraphs[0].font.color.rgb = STATUS_RGB.get(color_key, STATUS_RGB["yellow"])[1]


def _add_title_slide(prs: Presentation, company: str, date: str, total_respondents: int = 0) -> None:
    """Adiciona slide de capa com wallpaper, logo e textos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w = prs.slide_width
    h = prs.slide_height
    wallpaper_path = IMG_DIR / "wallpaper_PPT.png"
    logo_path = IMG_DIR / "fluir_logo.png"
    if wallpaper_path.exists():
        slide.shapes.add_picture(str(wallpaper_path), 0, 0, w, h)
    if logo_path.exists():
        logo_size = Inches(LOGO_SIZE_INCHES)
        logo_left = Inches(13.333 - LOGO_SIZE_INCHES)
        slide.shapes.add_picture(str(logo_path), logo_left, Inches(0.2), logo_size, logo_size)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    t = title_box.text_frame.paragraphs[0]
    t.text = "Fluir"
    t.font.size = Pt(44)
    t.font.bold = True
    t.font.color.rgb = RGBColor(0x0F, 0x4C, 0x75)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Relatorio de Diagnostico Psicossocial"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x32, 0x82, 0xB8)
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    resp_text = f" | {total_respondents} respondentes" if total_respondents else ""
    tb2.text_frame.paragraphs[0].text = f"Organizacao: {company}  |  Data: {date}{resp_text}"
    tb2.text_frame.paragraphs[0].font.size = Pt(14)


def _add_table_slide(
    prs: Presentation,
    title_text: str,
    headers: List[str],
    rows: List[List[Any]],
    status_col_idx: Optional[int] = None,
    status_map: Optional[Dict[str, str]] = None,
) -> None:
    """Adiciona slide com tabela (formato 16:9). Opcional: coluna de status com cores."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(12)
    title_box = slide.shapes.add_textbox(left, Inches(0.3), width, Inches(0.6))
    title_box.text_frame.paragraphs[0].text = title_text
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_height = Inches(0.35 * min(n_rows, 12))
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height).table
    status_map = status_map or {}
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        try:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_PRIMARY
        except Exception:
            pass
    for r_idx, row in enumerate(rows, 1):
        for c_idx, val in enumerate(row):
            if c_idx < n_cols:
                cell = table.cell(r_idx, c_idx)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                if status_col_idx is not None and c_idx == status_col_idx:
                    status_val = str(val).lower().strip()
                    status_key = "green" if "favor" in status_val else "yellow"
                    if "crit" in status_val or "critico" in status_val:
                        status_key = "red"
                    elif "aten" in status_val:
                        status_key = "yellow"
                    colors = STATUS_RGB.get(status_key, STATUS_RGB["yellow"])
                    try:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors[0]
                    except Exception:
                        pass


def export_pptx(
    survey: Dict[str, Any],
    respondents_data: List[Dict[str, Any]],
    dim_scores_agg: List[Dict[str, Any]],
    kpis: Dict[str, Any],
    summary: Dict[str, Any],
    recommendations: List[Dict[str, Any]],
    recommendations_prose: Dict[str, str],
) -> io.BytesIO:
    """Gera apresentacao PowerPoint executiva e retorna BytesIO.

    Ordem: capa, resumo executivo, KPIs cards, radar, barras, resumo numerico,
    dimensoes, respostas, recomendacoes, prosa, ROI, fechamento. Formato 16:9.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_16_9
    prs.slide_height = SLIDE_HEIGHT_16_9

    company = survey.get("company_name", "Empresa")
    date_str = datetime.now().strftime("%d/%m/%Y")
    total_resp = summary.get("total", 0) or len(respondents_data)

    # 1. Capa
    _add_title_slide(prs, company, date_str, total_resp)

    # 2. Resumo executivo (copywriting)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tit.text_frame.paragraphs[0].text = "Resumo Executivo"
    tit.text_frame.paragraphs[0].font.size = Pt(28)
    tit.text_frame.paragraphs[0].font.bold = True
    tit.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(2))
    body.text_frame.word_wrap = True
    body.text_frame.paragraphs[0].text = RESUMO_EXECUTIVO
    body.text_frame.paragraphs[0].font.size = Pt(16)
    body.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    meta = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(0.5))
    meta.text_frame.paragraphs[0].text = f"{company} | {total_resp} respondentes | {date_str}"
    meta.text_frame.paragraphs[0].font.size = Pt(12)
    meta.text_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY

    # 3. KPIs em cards
    _add_kpi_cards_slide(prs, kpis)

    # 4. Panorama por categoria (radar)
    if dim_scores_agg:
        cat_scores = _build_category_scores(dim_scores_agg)
        try:
            radar_buf = _render_radar_chart(cat_scores)
            radar_slide = prs.slides.add_slide(prs.slide_layouts[6])
            tbox = radar_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            tbox.text_frame.paragraphs[0].text = "Panorama por Categoria"
            tbox.text_frame.paragraphs[0].font.size = Pt(28)
            tbox.text_frame.paragraphs[0].font.bold = True
            tbox.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
            radar_slide.shapes.add_picture(radar_buf, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        except Exception:
            pass

    # 5. Comparativo por dimensao (barras)
    if dim_scores_agg:
        try:
            bar_buf = _render_bar_chart(dim_scores_agg)
            bar_slide = prs.slides.add_slide(prs.slide_layouts[6])
            tbox = bar_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            tbox.text_frame.paragraphs[0].text = "Comparativo por Dimensao"
            tbox.text_frame.paragraphs[0].font.size = Pt(28)
            tbox.text_frame.paragraphs[0].font.bold = True
            tbox.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
            bar_slide.shapes.add_picture(bar_buf, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        except Exception:
            pass

    # 6. Resumo numerico
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tbox.text_frame.paragraphs[0].text = "Resumo Geral"
    tbox.text_frame.paragraphs[0].font.size = Pt(28)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = f"{summary.get('green', 0)} dimensoes favoraveis | {summary.get('yellow', 0)} em atencao | {summary.get('red', 0)} criticas"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_PRIMARY

    # 7. Analise por dimensao (tabela com cores de status)
    dim_headers = ["#", "Dimensao", "Score", "Status", "Tipo", "Categoria", "Descricao"]
    dim_rows = [
        [idx, d["name"], f"{d['score']:.2f}", STATUS_LABEL.get(d["status"], ""),
         "Risco" if d["type"] == "risk" else "Recurso", d["category"], (d.get("description") or "")[:MAX_DIM_DESC_LEN]]
        for idx, d in enumerate(dim_scores_agg, 1)
    ]
    _add_table_slide(prs, "Analise por Dimensao", dim_headers, dim_rows, status_col_idx=3)

    # 8. Respostas individuais
    if respondents_data and DIMENSIONS:
        dim_ids = list(DIMENSIONS.keys())
        resp_headers = ["ID"] + [DIMENSIONS[d]["name"][:8] for d in dim_ids]
        resp_rows = []
        for r in respondents_data:
            row = [r.get("display_id", "")]
            for dim_id in dim_ids:
                s = r.get("scores", {}).get(dim_id)
                row.append(f"{s:.2f}" if s is not None else "-")
            resp_rows.append(row)
        if resp_rows:
            _add_table_slide(prs, "Respostas Individuais", resp_headers, resp_rows)

    # 9. Recomendacoes estruturadas
    if recommendations:
        rec_headers = ["#", "Prioridade", "Titulo", "Descricao"]
        rec_rows = [
            [idx, PRIORITY_LABEL.get(r.get("priority", ""), r.get("priority", "")), r.get("title", ""), (r.get("description") or "")[:MAX_REC_DESC_LEN]]
            for idx, r in enumerate(recommendations, 1)
        ]
        _add_table_slide(prs, "Recomendacoes Estruturadas", rec_headers, rec_rows)

    # 10. Introducao copywriting + Recomendacoes em prosa
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    intro_box.text_frame.word_wrap = True
    intro_box.text_frame.paragraphs[0].text = "Plano de Acao Recomendado"
    intro_box.text_frame.paragraphs[0].font.size = Pt(28)
    intro_box.text_frame.paragraphs[0].font.bold = True
    intro_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    intro_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.8))
    intro_text.text_frame.word_wrap = True
    intro_text.text_frame.paragraphs[0].text = INTRO_RECOMENDACOES
    intro_text.text_frame.paragraphs[0].font.size = Pt(12)
    intro_text.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    sections = [("Acoes imediatas", "imediata"), ("Curto prazo", "curto_prazo"), ("Medio prazo", "medio_prazo")]
    for title, key in sections:
        text = (recommendations_prose or {}).get(key) or ""
        if text.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tit_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
            tit_box.text_frame.paragraphs[0].text = title
            tit_box.text_frame.paragraphs[0].font.size = Pt(24)
            tit_box.text_frame.paragraphs[0].font.bold = True
            tit_box.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
            tx.text_frame.word_wrap = True
            tx.text_frame.paragraphs[0].text = text.strip()[:MAX_PROSE_LEN]
            tx.text_frame.paragraphs[0].font.size = Pt(12)

    # 11. Slide ROI/beneficios
    roi_slide = prs.slides.add_slide(prs.slide_layouts[6])
    roi_tit = roi_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    roi_tit.text_frame.paragraphs[0].text = ROI_TITULO
    roi_tit.text_frame.paragraphs[0].font.size = Pt(28)
    roi_tit.text_frame.paragraphs[0].font.bold = True
    roi_tit.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    roi_body = roi_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(4))
    tf = roi_body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(ROI_BULLETS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_PRIMARY
        if i > 0:
            p.space_before = Pt(8)

    # 12. Fechamento CTA
    footer_slide = prs.slides.add_slide(prs.slide_layouts[6])
    fb = footer_slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(2))
    fb.text_frame.word_wrap = True
    fb.text_frame.paragraphs[0].text = FECHAMENTO_CTA
    fb.text_frame.paragraphs[0].font.size = Pt(14)
    fb.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
